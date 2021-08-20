import numpy as np
import pandas as pd
from pandas import ExcelWriter
from pandas import ExcelFile
import xlrd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


#Importing Data using pandas and xlrd
df = pd.read_excel('2018_Q3_master_data.xlsx', index_col = 0)

#opening ppt file
path = 'kpi_template.pptx'
prs = Presentation(path)


def ex1(row_id, clist = []):

    #Choosing layout template for exhibit2 from the ppt file
    e1_layout = prs.slide_layouts[6]

    #adding a new slide to the ppt
    exhibit1 = prs.slides.add_slide(e1_layout)
    title = exhibit1.shapes.title
    title.text = "Food Attributes"

    #Finding how may editable fields exist.
    #Uncomment below 2 lines to find out
    #for shape in exhibit1.placeholders:
        #print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
    #Mapping textfields to variables
    text1 = exhibit1.placeholders[17]
    text2 = exhibit1.placeholders[18]
    text3 = exhibit1.placeholders[19]

    #Editing Subtexts
    text1.text = "Food Taste & Flavor in %"
    text2.text = "Food Quality in %"
    text3.text = "Food Quality When Ordered for Takeout*"
    
    #Comment adding to slide
    comment = exhibit1.placeholders[20]
    comment.text = "Total base: " +str(df.loc[row_id, "Base Size"])+" "+ "recent " +str(row_id) +" " +"guests"
    #Getting the data required for charts

###############################################################################################################
    for x in clist:

        if x == 16:
            ctype = "Food taste and flavor"
        elif x == 14:
            ctype = "Food quality"
        elif x == 15:
            ctype = "Food quality takeout"

        temp_list = [

            (row_id, round(df.loc[row_id, ctype]*100),1),
            (df.loc[row_id, "VA-1"], round(df.loc[df.loc[row_id, "VA-1"], ctype]*100),1),
            (df.loc[row_id, "VA-2"], round(df.loc[df.loc[row_id, "VA-2"], ctype]*100),1),
            (df.loc[row_id, "VA-3"], round(df.loc[df.loc[row_id, "VA-3"], ctype]*100),1),
            (df.loc[row_id, "VA-4"], round(df.loc[df.loc[row_id, "VA-4"], ctype]*100),1),
            (df.loc[row_id, "VA-5"], round(df.loc[df.loc[row_id, "VA-5"], ctype]*100),1),
            (df.loc[row_id, "VA-6"], round(df.loc[df.loc[row_id, "VA-6"], ctype]*100),1)
        
        ]
        temp_list.sort(reverse=False, key = lambda x: x[1])
    
    

        #Chart #1 code, Food taste & Flavor
        chart_data = ChartData()
        chart_data.categories = [
            temp_list[0][0],
            temp_list[1][0],
            temp_list[2][0],
            temp_list[3][0],
            temp_list[4][0],
            temp_list[5][0],
            temp_list[6][0],
        
        ]
        chart_data.add_series('Series 1', (
            temp_list[0][1],
            temp_list[1][1],
            temp_list[2][1],
            temp_list[3][1],
            temp_list[4][1],
            temp_list[5][1],
            temp_list[6][1],
        ))

    
        # add chart to slide --------------------
        create_chart1 = exhibit1.placeholders[x].insert_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, chart_data
        )

        
        #Getting chart from graphic frame
        chart1 = create_chart1.chart
        plot = chart1.plots[0]
        plot.has_data_labels = True

        #Editing grid lines
        chart1.value_axis.has_major_gridlines = False
        chart1.value_axis.has_minor_gridlines  =False



    

    #Tried to get the Chick-fil-A image. But, could not find it in the API
    #Will come back to this
    #logo = plot.categories[0].fill.type == MSO_FILL.PICTURE('logo/chickfila.png')

    
    
    
########################################################################
def ex2(row_id):

    #Choosing layout template for exhibit2 from the ppt file
    e2_layout = prs.slide_layouts[1]

    #adding a new slide to the ppt
    exhibit2 = prs.slides.add_slide(e2_layout)
    

    #editing 

    #Finding how may editable fields exist.
    #Uncomment below 2 lines to find out
    #for shape in exhibit2.placeholders:
        #print('%d %s' % (shape.placeholder_format.idx, shape.name))
    
    #Comment adding to slide
    comment = exhibit2.placeholders[26]
    comment.text = "Total base: " +str(df.loc[row_id, "Base Size"])+" "+ "recent " +str(row_id) +" " +"guests"
    #mapping textfields to variables
    title = exhibit2.shapes.title
    text1 = exhibit2.placeholders[20]
    text2 = exhibit2.placeholders[25]
    text3 = exhibit2.placeholders[24]
    text4 = exhibit2.placeholders[21]
    text5 = exhibit2.placeholders[22]
    text6 = exhibit2.placeholders[23]

    #Getting Competitors and the scores and 
    #sorting them according to descending percentage
    temp_List = [
        (df.loc[row_id, "VA-1"],round(df.loc[row_id, "VA1-score"]*100,1)), 
        (df.loc[row_id, "VA-2"],round(df.loc[row_id, "VA2-score"]*100,1)), 
        (df.loc[row_id, "VA-3"],round(df.loc[row_id, "VA3-score"]*100,1)), 
        (df.loc[row_id, "VA-4"],round(df.loc[row_id, "VA4-score"]*100,1)),
        (df.loc[row_id, "VA-5"],round(df.loc[row_id, "VA5-score"]*100,1)),
        (df.loc[row_id, "VA-6"],round(df.loc[row_id, "VA6-score"]*100,1))
        ]
    temp_List.sort(reverse=True, key = lambda x: x[1])
    

    #Image editing
    exhibit2.placeholders[14].insert_picture('logo/mcd.png')
    exhibit2.placeholders[19].insert_picture('logo/burgerking.png')
    exhibit2.placeholders[18].insert_picture('logo/kfc.png')
    exhibit2.placeholders[15].insert_picture('logo/subway.png')
    exhibit2.placeholders[17].insert_picture('logo/tacobell.png')
    exhibit2.placeholders[16].insert_picture('logo/wendys.png')
    #Subtext Editing
    #Getting Company name for the title
    title.text = "Top "+ str(row_id) + " Competitors"

    #Subtext of all the images editing. Text numbers are from left to right.
    text1.text =  str(temp_List[0][1]) +"%"+ " of recent "+ str(row_id) + " guests considered visiting " + temp_List[0][0]
    text2.text =  str(temp_List[1][1]) +"%"+ " of recent "+ " guests considered visiting " + temp_List[1][0]
    text3.text =  str(temp_List[2][1]) +"%"+ " of recent "+ " guests considered visiting " + temp_List[2][0]
    text4.text =  str(temp_List[3][1]) +"%"+ " of recent "+ " guests considered visiting " + temp_List[3][0]
    text5.text =  str(temp_List[4][1]) +"%"+ " of recent "+ " guests considered visiting " + temp_List[4][0]
    text6.text =  str(temp_List[5][1]) +"%"+ " of recent "+ " guests considered visiting " + temp_List[5][0]
    
    
################################################################################  

def main():
    #manually passing Chick-fil-A index from excel sheet
    #14,15,16 are the placeholders in ppt for charts
    ex1('Chick-fil-A',[14,15,16])
    ex2('Chick-fil-A')
    #Run these functions in the order you want the slides

main()
#Saving the new powerpoint file
prs.save('Chick-fil-A.pptx')
